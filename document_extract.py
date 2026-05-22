"""Extracción de texto desde PDF, Word, Excel y PowerPoint.

Devuelve siempre un dict {filename, content_type, text, char_count, error?}.
El texto se devuelve listo para pasar como contexto al LLM (limitado a MAX_CHARS).

Defensas activas:
- MAX_CHARS al output (límite de tokens al LLM).
- Zip-bomb protection en archivos basados en ZIP (DOCX/PPTX/XLSX): se
  inspecciona el archivo antes de pasarlo a las librerías, rechazando
  ratios de compresión sospechosos o tamaños descomprimidos enormes.
  Mitiga DoS por upload de archivo chico malicioso que se descomprime
  a varios GB y mata el worker uvicorn por OOM.
"""

from __future__ import annotations

import io
import zipfile
from pathlib import Path

MAX_CHARS = 60_000  # Tope de caracteres a devolver al LLM (≈ 15K tokens)

# ── Zip-bomb defense (M1) ────────────────────────────────────────────
# Un archivo Office (DOCX/PPTX/XLSX) bien formado típico:
#   - Tiene entre 5 y 50 entries internas.
#   - Pesa descomprimido entre 100 KB y 20 MB.
#   - Ratio de compresión: 2:1 a 10:1.
# Un zip-bomb intencional puede tener ratio >1000:1 y descomprimir a
# decenas de GB. Estos límites cortan los ataques sin afectar archivos
# legítimos grandes (un Excel con muchos datos puede llegar a 100 MB
# descomprimido en casos extremos — dejamos 200 MB de margen).
_MAX_UNCOMPRESSED_BYTES = 200 * 1024 * 1024  # 200 MB
_MAX_ZIP_ENTRIES = 5_000
_MAX_COMPRESSION_RATIO = 1_000  # Ratio total uncompressed/compressed


class ZipBombError(ValueError):
    """El archivo ZIP parece malicioso (zip-bomb)."""


def _check_zip_safety(data: bytes) -> None:
    """Inspecciona un ZIP sin descomprimirlo. Levanta ZipBombError si falla.

    Si los bytes no son un ZIP válido, retorna sin error (la función la
    llaman solo los extractores que esperan ZIP; si no lo es, el lib
    upstream va a fallar con un error específico de su formato).
    """
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            infos = zf.infolist()
            if len(infos) > _MAX_ZIP_ENTRIES:
                raise ZipBombError(
                    f"ZIP con demasiados archivos ({len(infos)} > {_MAX_ZIP_ENTRIES})"
                )
            total_uncompressed = sum(i.file_size for i in infos)
            if total_uncompressed > _MAX_UNCOMPRESSED_BYTES:
                mb = total_uncompressed // (1024 * 1024)
                raise ZipBombError(
                    f"ZIP se descomprime a {mb} MB "
                    f"(límite {_MAX_UNCOMPRESSED_BYTES // (1024 * 1024)} MB)"
                )
            total_compressed = sum(i.compress_size for i in infos)
            if total_compressed > 0:
                ratio = total_uncompressed / total_compressed
                if ratio > _MAX_COMPRESSION_RATIO:
                    raise ZipBombError(
                        f"Ratio de compresión sospechoso: "
                        f"{int(ratio)}:1 (límite {_MAX_COMPRESSION_RATIO}:1)"
                    )
    except zipfile.BadZipFile:
        # No es un ZIP válido — dejamos que el extractor upstream falle
        # con un mensaje más específico de su formato.
        return


def _truncate(text: str) -> str:
    text = text.strip()
    if len(text) <= MAX_CHARS:
        return text
    return text[:MAX_CHARS] + f"\n\n[...truncado en {MAX_CHARS} chars de {len(text)} totales]"


def _extract_pdf(data: bytes) -> str:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(data))
    out = []
    for i, page in enumerate(reader.pages, 1):
        try:
            txt = page.extract_text() or ""
        except Exception as e:
            txt = f"[Error extrayendo página {i}: {e}]"
        if txt.strip():
            out.append(f"--- Página {i} ---\n{txt.strip()}")
    return "\n\n".join(out) or "[PDF sin texto extraíble — posiblemente escaneado]"


def _extract_docx(data: bytes) -> str:
    _check_zip_safety(data)  # DOCX = ZIP → anti zip-bomb
    from docx import Document
    doc = Document(io.BytesIO(data))
    out = []
    for para in doc.paragraphs:
        if para.text.strip():
            out.append(para.text.strip())
    for ti, table in enumerate(doc.tables, 1):
        out.append(f"\n--- Tabla {ti} ---")
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            out.append(" | ".join(cells))
    return "\n".join(out) or "[Documento Word vacío]"


def _extract_pptx(data: bytes) -> str:
    _check_zip_safety(data)  # PPTX = ZIP → anti zip-bomb
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    out = []
    for si, slide in enumerate(prs.slides, 1):
        out.append(f"\n--- Slide {si} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(r.text for r in para.runs).strip()
                    if txt:
                        out.append(txt)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    out.append(" | ".join(cells))
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                out.append(f"[Notas]: {notes}")
    return "\n".join(out) or "[Presentación PowerPoint vacía]"


def _extract_xlsx(data: bytes) -> str:
    _check_zip_safety(data)  # XLSX = ZIP → anti zip-bomb
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    out = []
    for sheet in wb.worksheets:
        out.append(f"\n--- Hoja: {sheet.title} ---")
        rows_added = 0
        for row in sheet.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            if any(c.strip() for c in cells):
                out.append(" | ".join(cells))
                rows_added += 1
            if rows_added >= 500:
                out.append(f"[... hoja truncada en 500 filas]")
                break
    return "\n".join(out) or "[Excel vacío]"


def _extract_txt(data: bytes) -> str:
    for enc in ("utf-8", "latin-1", "cp1252"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


_EXTRACTORS = {
    ".pdf":  _extract_pdf,
    ".docx": _extract_docx,
    ".pptx": _extract_pptx,
    ".xlsx": _extract_xlsx,
    ".xlsm": _extract_xlsx,
    ".csv":  _extract_txt,
    ".txt":  _extract_txt,
    ".md":   _extract_txt,
}


def extract(filename: str, data: bytes) -> dict:
    """Extrae texto del archivo según extensión. Retorna dict con text + metadata."""
    ext = Path(filename).suffix.lower()
    fn = _EXTRACTORS.get(ext)
    if not fn:
        return {
            "filename": filename,
            "text": "",
            "char_count": 0,
            "error": f"Formato no soportado: {ext}. Aceptados: PDF, DOCX, PPTX, XLSX, CSV, TXT.",
        }
    try:
        text = fn(data)
        text = _truncate(text)
        return {
            "filename": filename,
            "text": text,
            "char_count": len(text),
        }
    except ZipBombError as e:
        # Mensaje específico — al usuario le interesa saber por qué se
        # rechazó, no sale por el bucket genérico.
        return {
            "filename": filename,
            "text": "",
            "char_count": 0,
            "error": f"Archivo rechazado por defensa anti-zip-bomb: {e}",
        }
    except Exception as e:
        return {
            "filename": filename,
            "text": "",
            "char_count": 0,
            "error": f"Error procesando archivo: {e}",
        }
